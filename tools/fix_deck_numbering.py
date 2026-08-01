"""
Fix the slide-number footers in PulseHR_Presentation.

Resolves P2-1. The footers are static text boxes, not real slide-number placeholders, so
they drifted as slides were inserted and reordered. In the 01-08-26 deck the sequence runs:

    1, 2, -, 4, 5, 6, 7, -, 9, 10, 11, 12, 11, 12, 13, 16, -, 16, 19, 20, -,
    22, 18, 19, 20, 21, -, 25, 24, -, 28, 32, 33

11, 12, 16, 19 and 20 each appear twice; it runs backwards at 22 -> 18; it jumps 13 -> 16.

This script finds each footer (a small text box near the bottom of the slide whose entire
content is an integer) and rewrites it to the slide's real position, preserving the run
formatting. Slides with no such box - the section dividers - are left alone.

    python tools/fix_deck_numbering.py <input.pptx> [output.pptx]

Writes to <input>_renumbered.pptx unless an output path is given. Never edits in place.
"""

from __future__ import annotations

import sys
from collections import Counter
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    sys.exit("python-pptx is required:  pip install python-pptx")

# EMU tolerance when matching a box against the deck's footer anchor (~0.05 cm).
POSITION_TOLERANCE = 20000


def digit_boxes(slide):
    """Every text box on the slide whose entire content is an integer."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not shape.text_frame.text.strip().isdigit():
            continue
        if shape.left is None or shape.top is None:
            continue
        yield shape


def find_footer_anchor(prs) -> tuple[int, int]:
    """
    Locate the footer by POSITION, not by "somewhere near the bottom".

    A naive "bottom 20% of the slide" rule is wrong: slide 25's step-by-step list puts
    step 7 at 83% of the slide height, and a percentage rule silently overwrites it with
    the page number. Real footers sit at one identical anchor on every slide, so the modal
    (left, top) across the whole deck identifies them unambiguously.
    """
    positions = Counter(
        (shape.left, shape.top) for slide in prs.slides for shape in digit_boxes(slide)
    )
    if not positions:
        sys.exit("no numeric text boxes found - nothing to renumber")

    (left, top), count = positions.most_common(1)[0]
    print(f"footer anchor: left={left} top={top}  (appears on {count} slides)\n")
    return left, top


def renumber(src: Path, dst: Path) -> None:
    prs = Presentation(str(src))
    anchor_left, anchor_top = find_footer_anchor(prs)

    changed, skipped, absent = 0, 0, []

    for index, slide in enumerate(prs.slides, start=1):
        boxes = [
            s
            for s in digit_boxes(slide)
            if abs(s.left - anchor_left) <= POSITION_TOLERANCE
            and abs(s.top - anchor_top) <= POSITION_TOLERANCE
        ]

        if not boxes:
            absent.append(index)
            continue

        for box in boxes:
            current = box.text_frame.text.strip()
            if current == str(index):
                skipped += 1
                continue

            # Rewrite the first run only, so font, size and colour survive.
            paragraph = box.text_frame.paragraphs[0]
            if paragraph.runs:
                paragraph.runs[0].text = str(index)
                for extra in paragraph.runs[1:]:
                    extra.text = ""
            else:
                box.text_frame.text = str(index)

            print(f"  slide {index:>2}: {current} -> {index}")
            changed += 1

    prs.save(str(dst))

    print(f"\n{changed} renumbered, {skipped} already correct")
    if absent:
        print(f"no number box on slides: {absent}")
        print("(section dividers - intentional, left alone)")
    print(f"\nwritten to {dst}")


def main() -> None:
    if len(sys.argv) < 2:
        sys.exit(__doc__)

    src = Path(sys.argv[1])
    if not src.exists():
        sys.exit(f"not found: {src}")

    dst = Path(sys.argv[2]) if len(sys.argv) > 2 else src.with_name(f"{src.stem}_renumbered.pptx")
    if dst.resolve() == src.resolve():
        sys.exit("refusing to overwrite the source deck - give a different output path")

    print(f"reading {src}\n")
    renumber(src, dst)


if __name__ == "__main__":
    main()
